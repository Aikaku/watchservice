"""
WatchService Agent 최종 발표 PPT 생성 스크립트
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK_BG   = RGBColor(0x0F, 0x17, 0x2A)   # #0f172a
SURFACE   = RGBColor(0x1E, 0x29, 0x3B)   # #1e293b
BLUE      = RGBColor(0x3B, 0x82, 0xF6)   # #3b82f6
RED       = RGBColor(0xEF, 0x44, 0x44)   # #ef4444
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)   # #f59e0b
GREEN     = RGBColor(0x22, 0xC5, 0x5E)   # #22c55e
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # #94a3b8
LIGHT     = RGBColor(0xF1, 0xF5, 0xF9)   # #f1f5f9

blank_layout = prs.slide_layouts[6]  # completely blank


def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=20, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, title, bullets, l, t, w, h,
                   title_color=BLUE, bullet_color=LIGHT,
                   bg_color=SURFACE):
    add_rect(slide, l, t, w, h, bg_color)
    add_text(slide, title, l+0.15, t+0.12, w-0.3, 0.4,
             size=14, bold=True, color=title_color)
    for i, b in enumerate(bullets):
        add_text(slide, f"• {b}", l+0.15, t+0.55+i*0.38, w-0.3, 0.4,
                 size=11, color=bullet_color)


# ──────────────────────────────────────────────
# Slide 1 — Title
# ──────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
add_bg(s1)

# top accent bar
add_rect(s1, 0, 0, 13.33, 0.06, BLUE)

# badge
add_rect(s1, 5.5, 1.4, 2.33, 0.38, RGBColor(0x1D, 0x4E, 0xD8))
add_text(s1, "SECURITY  ·  AI  ·  REAL-TIME", 5.5, 1.4, 2.33, 0.38,
         size=9, bold=True, color=RGBColor(0xBF, 0xDB, 0xFE),
         align=PP_ALIGN.CENTER)

# main title
add_text(s1, "WatchService Agent",
         1.5, 2.0, 10.33, 1.1,
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# subtitle blue
add_text(s1, "Real-time Ransomware Detection System",
         1.5, 3.1, 10.33, 0.7,
         size=24, bold=False, color=BLUE, align=PP_ALIGN.CENTER)

# divider
add_rect(s1, 5.2, 3.9, 2.93, 0.04, MUTED)

# description
add_text(s1,
         "Behavior-based AI detection  ·  Instant alert  ·  Gemini AI guidance",
         1.5, 4.1, 10.33, 0.5,
         size=14, color=MUTED, align=PP_ALIGN.CENTER)

# team info
add_text(s1, "Computer Science  |  2026",
         1.5, 5.8, 10.33, 0.5,
         size=12, color=MUTED, align=PP_ALIGN.CENTER)

# bottom bar
add_rect(s1, 0, 7.44, 13.33, 0.06, BLUE)


# ──────────────────────────────────────────────
# Slide 2 — System Overview & Problem
# ──────────────────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
add_bg(s2)
add_rect(s2, 0, 0, 13.33, 0.06, BLUE)
add_rect(s2, 0, 7.44, 13.33, 0.06, BLUE)

add_text(s2, "System Overview", 0.4, 0.15, 9, 0.6,
         size=28, bold=True, color=WHITE)
add_text(s2, "What problem does it solve?", 0.4, 0.72, 9, 0.4,
         size=14, color=MUTED)

# ── Problem box (left)
add_rect(s2, 0.35, 1.25, 3.9, 5.7, RGBColor(0x1C, 0x19, 0x17))
# red left border trick — thin rect
add_rect(s2, 0.35, 1.25, 0.07, 5.7, RED)
add_text(s2, "⚠  The Problem", 0.52, 1.35, 3.6, 0.45,
         size=14, bold=True, color=RED)

problems = [
    ("Ransomware encrypts files", "and demands ransom — data\nloss is irreversible"),
    ("Traditional AV misses", "new/unknown variants\n(signature-based limits)"),
    ("No real-time alert", "Users notice too late;\ndamage already done"),
]
for i, (h, b) in enumerate(problems):
    y = 1.9 + i * 1.6
    add_text(s2, h, 0.55, y, 3.5, 0.35, size=12, bold=True, color=LIGHT)
    add_text(s2, b, 0.55, y+0.33, 3.5, 0.6, size=10, color=MUTED)

# ── Arrow
add_text(s2, "→", 4.4, 3.7, 0.5, 0.5, size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ── Solution box (right)
add_rect(s2, 5.0, 1.25, 7.95, 5.7, SURFACE)
add_rect(s2, 5.0, 1.25, 0.07, 5.7, BLUE)
add_text(s2, "✓  Our Solution", 5.17, 1.35, 7.5, 0.45,
         size=14, bold=True, color=BLUE)

solutions = [
    ("Behavior-based Detection",
     "Monitors file activity patterns — not\nfile content. Works on any ransomware variant."),
    ("XGBoost AI Model",
     "9 behavioral features (read/write/delete/\nentropy change) → SAFE / WARNING / DANGER"),
    ("Instant Multi-channel Alert",
     "Fullscreen emergency popup + email\nnotification the moment DANGER is detected"),
    ("Gemini AI Guidance",
     "LLM generates step-by-step response\nguide tailored to the detected threat"),
]
for i, (h, b) in enumerate(solutions):
    y = 1.9 + i * 1.3
    add_rect(s2, 5.15, y, 7.6, 1.15, RGBColor(0x16, 0x23, 0x36))
    add_text(s2, h, 5.3, y+0.06, 7.2, 0.35, size=12, bold=True, color=BLUE)
    add_text(s2, b, 5.3, y+0.4, 7.2, 0.65, size=10, color=MUTED)


# ──────────────────────────────────────────────
# Slide 3 — Key Features
# ──────────────────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
add_bg(s3)
add_rect(s3, 0, 0, 13.33, 0.06, BLUE)
add_rect(s3, 0, 7.44, 13.33, 0.06, BLUE)

add_text(s3, "Key Features", 0.4, 0.15, 9, 0.6,
         size=28, bold=True, color=WHITE)
add_text(s3, "Core capabilities of WatchService Agent", 0.4, 0.72, 9, 0.4,
         size=14, color=MUTED)

features = [
    ("🔍", "Real-time File Monitoring",
     ["Java NIO WatchService",
      "Recursive folder watching",
      "CREATE / MODIFY / DELETE events",
      "3-second window aggregation"]),
    ("🤖", "AI-powered Detection",
     ["XGBoost binary classifier",
      "9 behavioral feature vectors",
      "Trained on 30,000 samples",
      "SAFE / WARNING / DANGER labels"]),
    ("📊", "Entropy Analysis",
     ["Multi-section sampling (40/40/20%)",
      "Encryption detection via entropy spike",
      "Noise filter: |diff| > 0.05 threshold",
      "file_encrypt_like_count feature"]),
    ("🚨", "Instant Alert System",
     ["Fullscreen emergency popup",
      "Email notification (SMTP)",
      "Gemini LLM response guidance",
      "Dashboard with history & stats"]),
    ("⚙️", "Persistent Settings",
     ["Fixed ownerKey 'default'",
      "Survives server restart",
      "Watch folders + exception rules",
      "Alert email configuration"]),
    ("📋", "Reporting & Admin",
     ["PDF report export",
      "CSV / JSON log download",
      "Admin dashboard",
      "False-positive reporting"]),
]

cols = 3
for idx, (icon, title, bullets) in enumerate(features):
    col = idx % cols
    row = idx // cols
    l = 0.35 + col * 4.3
    t = 1.3 + row * 2.85
    add_rect(s3, l, t, 4.1, 2.65, SURFACE)
    add_rect(s3, l, t, 4.1, 0.5, RGBColor(0x16, 0x23, 0x36))
    add_text(s3, f"{icon}  {title}", l+0.15, t+0.06, 3.8, 0.4,
             size=12, bold=True, color=BLUE)
    for i, b in enumerate(bullets):
        add_text(s3, f"• {b}", l+0.15, t+0.58+i*0.47, 3.8, 0.42,
                 size=10, color=LIGHT)


# ──────────────────────────────────────────────
# Slide 4 — Technical Architecture
# ──────────────────────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
add_bg(s4)
add_rect(s4, 0, 0, 13.33, 0.06, BLUE)
add_rect(s4, 0, 7.44, 13.33, 0.06, BLUE)

add_text(s4, "Technical Architecture", 0.4, 0.15, 12, 0.6,
         size=28, bold=True, color=WHITE)
add_text(s4, "System components and data flow", 0.4, 0.72, 9, 0.4,
         size=14, color=MUTED)

# ── Top row: 3 layers
layers = [
    ("Frontend  (React)", SURFACE, BLUE,
     ["React CRA  +  React Router",
      "Polling dashboard (5s interval)",
      "Notifications / Logs / Settings"]),
    ("Backend  (Spring Boot 3.5)", SURFACE, BLUE,
     ["Java 17  +  Gradle",
      "WatcherService → EventWindowAggregator",
      "SQLite (JDBC)  ·  JavaMailSender"]),
    ("AI Server  (Python)", SURFACE, BLUE,
     ["FastAPI  +  XGBoost",
      "POST /api/analyze  (9 features)",
      "Gemini API  (LLM guidance)"]),
]
for i, (title, bg, tc, items) in enumerate(layers):
    l = 0.35 + i * 4.3
    add_rect(s4, l, 1.25, 4.1, 2.5, bg)
    add_rect(s4, l, 1.25, 4.1, 0.45, RGBColor(0x16, 0x23, 0x36))
    add_text(s4, title, l+0.15, 1.3, 3.8, 0.38, size=12, bold=True, color=tc)
    for j, item in enumerate(items):
        add_text(s4, f"• {item}", l+0.15, 1.8+j*0.47, 3.8, 0.42,
                 size=10, color=LIGHT)
    if i < 2:
        add_text(s4, "↔", 4.45+i*4.3, 2.35, 0.3, 0.4,
                 size=20, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

# ── Data flow
add_text(s4, "Data Flow", 0.4, 3.95, 4, 0.4, size=13, bold=True, color=MUTED)
flow_steps = [
    ("1", "File Event", "NIO WatchService detects CREATE / MODIFY / DELETE", BLUE),
    ("2", "3s Window", "Aggregate events → compute 9 feature vectors", BLUE),
    ("3", "AI Analyze", "XGBoost returns score → SAFE / WARNING / DANGER", AMBER),
    ("4", "DANGER?", "Fullscreen popup + Email + Gemini guidance", RED),
]
for i, (num, step, desc, col) in enumerate(flow_steps):
    l = 0.35 + i * 3.2
    add_rect(s4, l, 4.45, 3.0, 2.7, SURFACE)
    add_rect(s4, l, 4.45, 0.45, 2.7, col)
    add_text(s4, num, l+0.08, 4.55, 0.3, 0.4, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s4, step, l+0.6, 4.52, 2.3, 0.4, size=12, bold=True, color=WHITE)
    add_text(s4, desc, l+0.55, 4.98, 2.35, 1.8, size=10, color=MUTED, wrap=True)
    if i < 3:
        add_text(s4, "→", 3.5+i*3.2, 5.55, 0.2, 0.4,
                 size=16, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# Slide 5 — Demo Scenario
# ──────────────────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
add_bg(s5)
add_rect(s5, 0, 0, 13.33, 0.06, RED)
add_rect(s5, 0, 7.44, 13.33, 0.06, RED)

add_text(s5, "Demo", 0.4, 0.15, 9, 0.6,
         size=28, bold=True, color=WHITE)
add_text(s5, "Live demonstration scenario", 0.4, 0.72, 9, 0.4,
         size=14, color=MUTED)

steps = [
    ("Step 1", "Launch & Setup",
     "Open WatchService Agent  →  Start monitoring\nSettings → Add watch folder (e.g. Desktop/test)",
     BLUE),
    ("Step 2", "Configure Email Alert",
     "Settings → Alert Email → Enter recipient\nSend test email → Verify delivery in inbox",
     BLUE),
    ("Step 3", "Simulate Normal Activity",
     "Create / edit / delete a few text files\nDashboard shows SAFE (green) — no alert",
     GREEN),
    ("Step 4", "Simulate Ransomware Behavior",
     "Run simulate script: rapid mass rename +\nhigh-entropy overwrite on 20+ files",
     AMBER),
    ("Step 5", "DANGER Detection",
     "Fullscreen emergency popup appears\nEmail alert delivered to inbox",
     RED),
    ("Step 6", "Response & Review",
     "Gemini AI guidance displayed on screen\nDashboard → Notification detail → history log",
     BLUE),
]

for i, (tag, title, desc, col) in enumerate(steps):
    col_i = i % 2
    row_i = i // 2
    l = 0.35 + col_i * 6.5
    t = 1.3 + row_i * 1.95
    add_rect(s5, l, t, 6.2, 1.8, SURFACE)
    add_rect(s5, l, t, 0.07, 1.8, col)
    add_text(s5, tag, l+0.2, t+0.1, 1.2, 0.35, size=9, bold=True, color=col)
    add_text(s5, title, l+0.2, t+0.42, 5.8, 0.4, size=13, bold=True, color=WHITE)
    add_text(s5, desc, l+0.2, t+0.85, 5.8, 0.85, size=10, color=MUTED)


# ── Save
out = "/Users/sanghyeok/Desktop/Project/통합/문서/WatchService_Agent_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
